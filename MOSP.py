#################
### LIBRARIES ###
#################

# Import Python Libraries 
import pandas as pd 
import numpy as np

# Import Visualization Libraries 
import matplotlib.pyplot as plt
import seaborn as sns

# Import Powerpoint Libraries 
from pptx import Presentation 
from pptx.util import Inches, Pt 
from pptx.dml.color import RGBColor 
from pptx.chart.data import CategoryChartData
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN

###################
### IMPORT DATA ###
###################

# Import Data

def import_data(file):
    def import_sheet(sheet_name):
        data = pd.read_excel(file, sheet_name=sheet_name)
        data = data.fillna(' ')
        return data
    
    sheets = ['AppDev', 'AppMon', 'PerDev', 'PerMon','PSIDev','PSIMon','KSGini']
    dfs = {}
    
    for sheet_name in sheets:
        dfs[sheet_name] = import_sheet(sheet_name)
    
    return dfs['AppDev'], dfs['AppMon'], dfs['PerDev'], dfs['PerMon'], dfs['PSIDev'], dfs['PSIMon'], dfs['KSGini']


### CREATE POWERPOINT SLIDE ###

def pptx(file_name):  
    
    global slide 
    global prs 
    
    # Load the pre-existing PowerPoint file 
    prs = Presentation(file_name)  

    # Set PowerPoint layout 
    slide = prs.slides.add_slide(prs.slide_layouts[3]) 
    
    return slide 

###############################
### CREATE POWERPOINT TABLE ### 
###############################

def tableDF(df, x, y, widthOne, widthTwo, height, prsTitle, prsSubheading, prsText, model, fontC, fontD, fontH): 
    # Add the dataframe to the PowerPoint slide as a table
    rows, cols = df.shape
    shape = slide.shapes.add_table(rows + 2, cols, Inches(x), Inches(y), width=Inches(widthOne), height=Inches(height)) 
    
    table = shape.table   
    
    # Adjust the column widths 
    table.columns[0].width = Inches(widthOne)  # Set the width of the first column to 2 inches
    for col in list(table.columns)[1:]: 
        col.width = Inches(widthTwo)  # Set the width of each remaining column to 0.5 inches
        
    # Set the style ID for the table
    tbl = shape._element.graphic.graphicData.tbl
    style_id = '{616DA210-FB5B-4158-B5E0-FEB733F419BA}'
    tbl[0][-1].text = style_id
    
    # Merge cells in the first row and add "development" text
    cell_range = table.cell(0, 0)
    cell_range.merge(table.cell(0, cols - 1))
    cell_range.text = model
    cell_range.text_frame.paragraphs[0].font.size = Pt(fontH)
    cell_range.text_frame.paragraphs[0].font.name = 'Calibri' 
    cell_range.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Copy content from original table to the new table
    for i in range(rows + 1):
        for j in range(cols):
            cell = table.cell(i, j)
            new_cell = table.cell(i + 1, j)
            new_cell.text = cell.text
            new_cell.text_frame.paragraphs[0].font.size = cell.text_frame.paragraphs[0].font.size
            new_cell.text_frame.paragraphs[0].font.name = cell.text_frame.paragraphs[0].font.name
    
    # Add the column names and set font size and style
    for i, col_name in enumerate(df.columns):
        cell = table.cell(1, i)
        cell.text = col_name
        cell.text_frame.paragraphs[0].font.size = Pt(fontC)
        cell.text_frame.paragraphs[0].font.name = 'Calibri'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add the data to the table and set font size and style
    for i, row in df.iterrows():
        for j, val in enumerate(row):
            cell = table.cell(i + 2, j)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(fontD)
            cell.text_frame.paragraphs[0].font.name = 'Calibri'
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
    # Set the fill color of the cells in the first row to dark blue 
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(55, 86, 119)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
    # Set the fill color of the cells in the second row to dark blue
    for cell in table.rows[1].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(55, 86, 119)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Set the fill color of the cells in the data rows to white
    for i in range(2, len(table.rows)):
        for cell in table.rows[i].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Table Margin 
    for row in table.rows:
        for cell in row.cells:
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            
    # Add a title to the slide
    title_shape = slide.shapes.title
    title_shape.text = prsTitle
    title_shape.text_frame.paragraphs[0].font.name = 'Bebas Neue' 
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Add a subheading to the slide
    sub_title = slide.placeholders[1]
    sub_title.text = prsSubheading  
    sub_title.text_frame.paragraphs[0].font.size = Pt(20)

    # Add text to the slide
    body_shape = slide.shapes[2].text_frame 
    tf = body_shape 
    tf.text = prsText
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    return table 
    

##################################
### POWER POINT VISUALIZATIONS ###
##################################

# VISUAL with one plot 
    
def visual(plotType, df, x1, y1, xAxisTitle, yAxisTitle, graphTitle, yGraphPost): 
   
    # Figure Settings 
    fig, ax = plt.subplots(figsize=(14, 6)) 
        
    # Plot Settings 
    if plotType == 'line': 
        sns.lineplot(data=df, x = x1, y = y1, ax = ax, marker='o', color='#375677')
    if plotType == 'bar': 
        sns.barplot(data=df, x=x1, y=y1, ax=ax, color='#375677')
        
    # Graph Settings 
    plt.title(graphTitle, fontsize= 18) 
    plt.xlabel(xAxisTitle, fontsize= 16) 
    plt.ylabel(yAxisTitle, fontsize= 14)        
    plt.xticks(fontsize = 14) 
    plt.yticks(fontsize = 14) 
    plt.rcParams['font.family'] = 'Calibri' 
        
    # Set background color to white
    ax.set_facecolor('white')
        
    # Set color of axis lines
    ax.spines['bottom'].set_color('black')
    ax.spines['left'].set_color('black')
        
    # Save the plot in powerpoint 
    plt.savefig(graphTitle + '.png')
    pic = slide.shapes.add_picture(graphTitle + '.png', Inches(7), Inches(yGraphPost), width = Inches(6.5), height = Inches(3))
    
    plt.close('all') 
    
# VISUAL with two plots 
     
def visualTwo(plotType, df1, df2, x1, x2, y1, y2, graphTitle, xAxisTitle, yAxisTitle, labelOne, labelTwo, yGraphPost): 
    
    # Figure Settings 
    fig, ax = plt.subplots(figsize=(14, 6)) 
    
    # Plot Settings                                                               
    if plotType == 'lineTwo':                                            
        sns.lineplot(data = df1, x=x1, y=y1, ax = ax, label=labelOne, marker='o', color='#375677')
        sns.lineplot(data = df2, x=x2, y=y2, ax = ax, label=labelTwo, marker='o', color='#8B0000')
    if plotType == "barLine": 
        sns.barplot(data = df1, x=x1, y= y1, ax = ax, label= labelOne, color='#375677')
        sns.lineplot(data = df2, x=x2, y= y2, ax = ax, label= labelTwo, marker='o', color='#8B0000')
    
    #Graph Settings 
    plt.title(graphTitle, fontsize= 18) 
    plt.xlabel(xAxisTitle, fontsize= 16) 
    plt.ylabel(yAxisTitle, fontsize= 16)        
    plt.xticks(fontsize = 14) 
    plt.yticks(fontsize = 14) 
    plt.legend(fontsize = 14)  
    plt.rcParams['font.family'] = 'Calibri'

    
    # Set background color to white
    ax.set_facecolor('white')
        
    # Set color of axis lines
    ax.spines['bottom'].set_color('black') 
    ax.spines['left'].set_color('black')  
         
    # Save the plot in powerpoint 
    plt.savefig(graphTitle + '.png')
    pic = slide.shapes.add_picture(graphTitle + '.png', Inches(7), Inches(yGraphPost), width = Inches(6.5), height = Inches(3))

    plt.close('all')
    
# VISUAL EMPTY 

def visualThree(graphTitle, xAxisTitle, yAxisTitle, yGraphPost): 
    
    # Figure Settings 
    fig, ax = plt.subplots(figsize=(14, 6)) 

    #Graph Settings 
    plt.title(graphTitle, fontsize= 18) 
    plt.xlabel(xAxisTitle, fontsize= 16) 
    plt.ylabel(yAxisTitle, fontsize= 16)        
    plt.xticks(fontsize = 14) 
    plt.yticks(fontsize = 14)   
    plt.rcParams['font.family'] = 'Calibri'
    
    # Save the plot in powerpoint 
    plt.savefig(graphTitle + '.png')
    pic = slide.shapes.add_picture(graphTitle + '.png', Inches(7), Inches(yGraphPost), width = Inches(6.5), height = Inches(3))

    plt.close('all') 

#################################
### SAVE POWER POINT CONTENTS ###
################################# 

def save(file_name):
    
    # Save the updated PowerPoint file 
    prs.save(file_name)

##############
### SLIDES ### 
##############

### FIRST SLIDE: POPULATION STABILITY INDEX ###  

# Calculate No. of Applications, Distribution % 
def calculate_distribution(data, bins):  
        counts = []  
        total_customers = len(data['applicationScore'])   
        
        for bin_range in bins:   
            lower, upper = map(int, bin_range.split('-'))  
            accounts_in_bin = data[(data['applicationScore'] >= lower) & (data['applicationScore'] <= upper)]  
            bin_count = len(accounts_in_bin)  
            counts.append(bin_count)  

        distribution = np.divide(counts, total_customers) * 100
        counts = pd.Series(counts)
        counts = counts.round(0).astype(int)

        return pd.DataFrame({'Score Band': bins, 'No. of Applications': counts, 'Distribution %': distribution.round(2)})  
        
    
### SECOND SLIDE: Predictive POWER - KS and GINI ### 

#Create Good and Bad Table  
def calculate_perf(df1, df2, months, delinquency):  
    df = pd.merge(df1, df2, on="accountNumber", how="inner")

    # convert the reportingDate and applicationDate columns to datetime objects
    df["reportingDate"] = pd.to_datetime(df["reportingDate"]) 
    df["applicationDate"] = pd.to_datetime(df["applicationDate"]) 

    #Offset applicationDate column  
    df["applicationDate"] = df["applicationDate"] + pd.offsets.MonthEnd(0)

    # Extract the date part of the reportingDate column 
    df["reportingDate"] = pd.to_datetime(df["reportingDate"]).dt.date

    # Extract the date part of the applicationDate column 
    df["applicationDate"] = pd.to_datetime(df["applicationDate"]).dt.date

    # create a new column called "time" that calculates the difference in months between reportingDate and applicationDate
    df["time"] = ((df["reportingDate"] - df["applicationDate"]) / np.timedelta64(1, 'M')).astype(int)

    # filter the goodBad to only include rows where the time is less than or equal to the specified months
    df = df[df["time"] <= months].copy()
    df.loc[:, "badFlag"] = np.where(df["delinquency"] >= delinquency, 1, 0) 

    return df 

# Calculate the number / distribution / cumulative of good and bad customers in each bin
def calculate_goodBad(data, flag, df1, df2):      
    counts = [] 
    
    df = pd.merge(df1, df2, on='accountNumber')
    
    num_bins = 10
    df['ScoreBucket'] = pd.qcut(df['applicationScore'], q=num_bins, labels=False)

    # Define bin edges and labels
    bins = pd.IntervalIndex.from_breaks(df['applicationScore'].quantile([i/num_bins for i in range(num_bins + 1)]))
    bin_labels = [f"{int(b.left)} - {int(b.right)}" for b in bins]

    for bin in bin_labels:
        lower, upper = map(int, bin.split('-')) 
        accounts_in_bin = data[(data['badFlag'] == flag) & (data['applicationScore'] >= lower) & (data['applicationScore'] <= upper)]
        bin_count = len(accounts_in_bin) 
        counts.append(bin_count) 

    counts = np.array(counts, dtype=float)  
    total = np.sum(data['badFlag'] == flag) 
    distribution = ((counts / total) * 100).round(2) 
    cumulative = (np.cumsum(distribution)).round(2) 

    return counts, total, distribution, cumulative  

### THIRD SLIDE: CHARACTER STABILITY INDEX ###

# Calculate / Group Distribution and Volume
def calculate_char(data, column):  
     
            score_column = 'wtVar' + column[3:].zfill(3)
            df = data.groupby([column, score_column])['applicationReference'].count().rename('Volume').reset_index()
            df['Distribution'] = ((df['Volume'] / df['Volume'].sum()) * 100).round(2)
            
            df.rename(columns={column: column, score_column: 'Score', 'Distribution': 'Dist.','Volume': 'Volume'}, inplace=True) 
            
            return df         

#Create GB table 
def calculate_char_two(data, column):  

        df1 = data[data['badFlag'] == 1].groupby([column]) ["badFlag"].count().rename('No. of Bads').reset_index()
        df2 = data[data['badFlag'] == 0].groupby([column])["badFlag"].count().rename('No. of Goods').reset_index()

        # compute distribution of bads and goods 

        totalBads =  df1['No. of Bads'].sum()
        totalGoods = df2['No. of Goods'].sum()

        df1['badDist'] = (df1['No. of Bads'] / totalBads *100).round(2)
        df2['goodDist'] = (df2['No. of Goods'] / totalGoods *100).round(2)

        GB = pd.merge(df1, df2, on = column, how = 'outer').fillna(0) 

        GB['Bad Rate'] = ((GB['No. of Bads'] / (GB['No. of Bads'] + GB['No. of Goods']))*100).round(2)

        GB['WOE'] = (np.log(GB['goodDist'] / GB['badDist'])).round(2)

        GB['IV'] = ((GB['goodDist'] - GB['badDist']) * GB['WOE']).round(2) 

        return GB  
    
##################
### EQUATIONS ####
##################

# PSI 
def psi(dev, mon): 
    
    # Convert 'Distribution %' columns to numeric values
    dev['Distribution %'] = pd.to_numeric(dev['Distribution %'],errors='coerce')
    mon['Distribution %'] = pd.to_numeric(mon['Distribution %'],errors='coerce') 
    
    values = ((dev['Distribution %'] - mon['Distribution %']) * np.log(dev['Distribution %'] / mon['Distribution %'])).round(2)
    
    sum = (np.sum(values)).round(2)
    
    return values, sum 


########################
### CLASSIFICATIONS ####
########################

# Classify PSI 
def psi_classify(value, df):  
    
    if df.empty: 
        psiValue = pd.DataFrame({'Value':  [" "], "Classification Quality" :  [" "]})
    else: 
        if value < 0.1: 
            classification = 'Satisfactory'                                   
        elif 0.1 <= value <= 0.25:                 
            classification = 'Warning'
        else:                                  
            classification = "Population Shift"                                               

        psiValue = pd.DataFrame({'Value': value, "Classification Quality" : [classification]})

    return psiValue   

# Classify KS 
def ks_classify(value, df):  
    
    if df.empty: 
        ksValue = pd.DataFrame({'Value':  [" "], "Classification Quality" :  [" "]})
    else: 
        if value < 0.15:
            classification = 'Poor Fit'
        elif 0.15 <= value < 0.25:
            classification = 'Average Fit'
        elif 0.25 <= value <= 0.45:
            classification = 'Good Fit'
        else:
            classification = "Strong Fit"
        
        ksValue = pd.DataFrame({'Value': value, "Classification Quality" : [classification]})

    return ksValue 

# Classify Gini 
def gini_classify(value, df):  
    if df.empty: 
        giniValue = pd.DataFrame({'Value':  [" "], "Classification Quality" :  [" "]})
    else: 
        if value < 0.25:     
            classification = 'Poor Fit'
        elif 0.25 <= value < 0.35:   
            classification = 'Average Fit'
        elif 0.35 <= value <= 0.55:  
            classification = 'Good Fit'
        else:
            classification = "Strong Fit"
                                            

        giniValue = pd.DataFrame({'Value': value, "Classification Quality" : [classification]})
        
    return giniValue 

#Classify CSI 
def csi_classify(value, df): 
    if df.empty: 
        csiValue = pd.DataFrame({'Value':  [" "], "Classification Quality" :  [" "]})
    else: 
        if value < 0.1:
            classification = 'Insignificant change'
        elif 0.1 <= value <= 0.25:
            classification = 'Some minor change'
        else:
            classification = "Major shift in population"
                                        
        csiValue = pd.DataFrame({'Value': value, "Classification Quality" : [classification]})
        
    return csiValue 

#Classify IV 
def iv_classify(value, df): 
    if df.empty: 
        ivValue = pd.DataFrame({'Value':  [" "], "Classification Quality" :  [" "]})
    else: 
        if value < 0.02:
            classification = 'Poor'
        elif 0.02 <= value < 0.1:
            classification = 'Average'
        elif 0.1 <= value <= 0.3:
            classification = 'Good'
        else:
            classification = "Strong"
                                            
        ivValue = pd.DataFrame({'Value': value, "Classification Quality" : [classification]})
        
    return ivValue 

##############################################################
### TABLE FORMATTING [PSI, IV, KS, GINI, CSI]: ###############
##############################################################

# Table formatting 
def apply_formatting(table):
    last_row = table.rows[len(table.rows) - 1]  # Get the last row

    for cell in last_row.cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(200, 200, 200) # Set the background grey 

        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True  # Set the font to bold
                
# Classifications 
def apply_formatting2(table, type, metric): 
    last_row = table.rows[len(table.rows) - 1]  # Get the last row
    last_cell = last_row.cells[len(last_row.cells) - 1]  # Get the last cell in the last row
    
    if type == "psi":
        if metric < 0.1:
            classification = RGBColor(144, 238, 144)  # Light Green
        elif 0.1 <= metric <= 0.25:
            classification = RGBColor(255, 204, 153)  # Light Orange
        else:
            classification = RGBColor(255, 153, 153)  # Light Red
            
    elif type == "ks":
        if metric < 0.15:
             classification = RGBColor(255, 153, 153)  # Light Red
        elif 0.15 <= metric < 0.25:
            classification = RGBColor(255, 204, 153)  # Light Orange
        elif 0.25 <= metric <= 0.45:
            classification = RGBColor(144, 238, 144)  # Light Green
        else:
            classification = RGBColor(152, 251, 152)  # Dark Green
            
    elif type == "gini":
        if metric < 0.25:
             classification = RGBColor(255, 153, 153)  # Light Red
        elif 0.25 <= metric < 0.35:
            classification = RGBColor(255, 204, 153)  # Light Orange
        elif 0.35 <= metric <= 0.55:
            classification = RGBColor(144, 238, 144)  # Light Green
        else:
            classification = RGBColor(152, 251, 152)  # Dark Green
            
    elif type == "csi":
        if metric < 0.1:
             classification = RGBColor(255, 153, 153)  # Light Red
        elif 0.1 <= metric < 0.25:
            classification = RGBColor(255, 204, 153)  # Light Orange
        else: 
            classification = RGBColor(144, 238, 144)  # Light Green
            
    elif type == "iv":
        if metric < 0.02:
             classification = RGBColor(255, 153, 153)  # Light Red
        elif 0.02 <= metric < 0.1:
            classification = RGBColor(255, 204, 153)  # Light Orange
        elif 0.1 <= metric <= 0.3:
            classification = RGBColor(144, 238, 144)  # Light Green
        else:
            classification = RGBColor(152, 251, 152)  # Dark Green

    for paragraph in last_cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = True  # Set the font to bold
            
############################            
### SCORECARD MONITORING ###
############################
def scorecard(file, bins, months, delinquency, name_dict, file_name): 
    
### Load Data ### 
    AppDev, AppMon, PerDev, PerMon, PSIDev, PSIMon, KSGini = import_data(file) 

###############################################
### FIRST SLIDE: POPULATION STABILITY INDEX ### 
###############################################

    # Overall Distribution Development Data 
    if AppDev.empty:  
        overallDistDev = pd.DataFrame(PSIDev)
        overallDistDev = overallDistDev.fillna(' ') 
            
    else: 
        overallDistDev = calculate_distribution(AppDev, bins)        
        
        # Create a DataFrame for the total row           
        total_row1 = pd.DataFrame({'Score Band': ['Total'],  
                               'No. of Applications': [np.sum(overallDistDev['No. of Applications'])], 
                               'Distribution %': [np.sum(overallDistDev['Distribution %']).round(2)]
                              })  
        
        overallDistDev = overallDistDev.append(total_row1, ignore_index=True) 
    
    # Overall Distribution Monitoring Data 
    if AppMon.empty: 
        overallDistMon = pd.DataFrame(PSIMon) 
        overallDistMon = overallDistMon.fillna(' ') 
        
        psi_values, psi_sum = psi(overallDistDev, overallDistMon)
        
        psi_summary = psi_classify(psi_sum, AppMon)
        
    else: 
        overallDistMon = calculate_distribution(AppMon, bins)  

        # Calculate the PSI value for each bin
        psi_values, psi_sum = psi(overallDistDev, overallDistMon)

        # Add the PSI column to overallDistMon  
        overallDistMon['PSI'] = psi_values 
    
        psi_summary = psi_classify(psi_sum, AppMon)
        
        # Create a DataFrame for the total row
        total_row2 = pd.DataFrame({'No. of Applications': [np.sum(overallDistMon['No. of Applications'])],
                               'Distribution %': [np.sum(overallDistMon['Distribution %']).round(2)],
                               'PSI': [np.sum(psi_values).round(2)]
                              })
        
        overallDistMon = overallDistMon.iloc[:, 1:]
        
        overallDistMon = overallDistMon.append(total_row2, ignore_index=True) 
    
    ### POWER POINT ###              

    # TABLE SETTINGS:                      
    Width1 = 1
    Width2 = 1.5
    Height = 1
    yHeight = 0.25 
    Font = 10 

    pptx(file_name)  

    #TABLE OVERAL DIST. DISTRIBUTION TABLE       
    table = tableDF(overallDistDev, 0.4, 3.1, 1, 1, 3.5, "Population Stability Index", "Personal Loan", " ", 'Development', 10, 10, 10) 
    apply_formatting(table) 

    # TABLE OVERAL DIST. MONITORING TABLE      
    table = tableDF(overallDistMon, 3.8, 3.1, 1, 1, 3.5, "Population Stability Index", "Personal Loan", " ", 'Monitoring', 10, 10, 10) 
    apply_formatting(table) 

    # TABLE PSI CLASSIFICATION TABLE     
    table3 = tableDF(psi_summary, 10.7, yHeight, Width1, Width2, Height, "Population Stability Index", "Personal Loan", " ", 'PSI', Font, Font, Font) 

    if not AppMon.empty:  
        apply_formatting2(table3, 'psi', psi_sum) 
    else: 
        pass 
    
    # VISUAL ONE           
    if not AppMon.empty: 
        visual("line", overallDistMon,  overallDistDev['Score Band'], 'PSI', 'Score Band', 'PSI', 'Population Stability Index', 1.57)  
    else: 
        visualThree('Population Stability Index','ScoreBand','PSI', 1.57) 
        
    # VISUAL TWO  
    if not AppDev.empty and not AppMon.empty: 
        visualTwo("lineTwo", overallDistDev, overallDistMon, overallDistDev["Score Band"], overallDistDev["Score Band"], 'No. of Applications', 'No. of Applications', 'Monitoring vs. Development Distribution', 'Score Band', 'Distribution %','Monitoring Distribution %', 'Development Distribution  %', 4.38 )
    else: 
        pass 
    
    if AppMon.empty and not AppDev.empty: 
        visual("line", overallDistDev, 'Score Band', 'No. of Applications', 'Score Band', 'Distribution %', 'Development Distribution', 4.38) 
    else: 
        pass 
    
    if AppMon.empty and AppDev.empty: 
        visualThree('Monitoring vs. Development Distribution','ScoreBand','Distribution %', 4.38) 
    else: 
        pass 
    
    # SAVE FILE 
    save(file_name) 
    
####################################################
### SECOND SLIDE: Predictive Power - KS and Gini###
###################################################

# KS GINI  
 
    if PerMon.empty: 
        ksGini = pd.DataFrame(KSGini) 
        
        ksTotal = 0  
        giniTotal = 0  
        
        ks_summary = ks_classify(ksTotal, AppMon)
        gini_summary = gini_classify(giniTotal, AppMon)
        
    else: 
        
        # Binning 
        df = pd.merge(AppMon, PerMon, on='accountNumber')
    
        num_bins = 10
        df['ScoreBucket'] = pd.qcut(df['applicationScore'], q=num_bins, labels=False)

        # Define bin edges and labels
        bins = pd.IntervalIndex.from_breaks(df['applicationScore'].quantile([i/num_bins for i in range(num_bins + 1)]))
        bin_labels = [f"{int(b.left)} - {int(b.right)}" for b in bins]

        # Intialize tables 

        tablesDev = {} 
        tablesMon = {}
        var_columns = AppDev.columns[AppDev.columns.str.startswith("var")]

        # Calculate Performance Information 
        goodbadDev = calculate_perf(AppDev, PerMon, months, delinquency)  
        goodBadMon = calculate_perf(AppMon, PerMon, months, delinquency)    

        Bads, totalBads, distributionBad, cumulativeBad = calculate_goodBad(goodBadMon, 1, AppMon, PerMon)    
        Goods, totalGoods, distributionGood, cumulativeGood = calculate_goodBad(goodBadMon, 0, AppMon, PerMon)  

        goodBadRatio = (Goods / Bads).round(2) 

        badRate = ((Bads / (Bads + Goods)) * 100).round(2) 

        # Calculate KS 

        KS = (abs(cumulativeBad - cumulativeGood)).round(2) 

        # Add 0 as the first element 
        cumulativeGood = np.append(cumulativeGood, 0)

        # Calculate the differences between consecutive elements 
        cumulative_diff = np.diff(cumulativeGood)       

        # Calculate the cumulative sum 
        cumulative_sum = np.cumsum(cumulativeBad) 

        gini = ((cumulative_diff) * ((cumulative_sum) / 2)).round(2)     

        # Create a dataframe to store the result 
        ksGini = pd.DataFrame({'Score Band': bin_labels}) 
        ksGini['No. of Bads'] = Bads 
        ksGini['No. of Goods'] = Goods
        ksGini['% Bads'] = distributionBad 
        ksGini['% Goods'] = distributionGood
        ksGini['G:B Odds'] = goodBadRatio
        ksGini['Bad Rate'] = badRate
        ksGini['KS'] = KS
        ksGini['Gini'] = gini

        # Create a DataFrame for the total row 
        total_row = pd.DataFrame({'Score Band': ['Total'],
                                  'No. of Bads': np.sum(Bads),
                                  'No. of Goods': np.sum(Goods),
                                  '% Bads': np.sum(distributionBad).round(2),
                                  '% Goods': np.sum(distributionGood).round(2),
                                  'G:B Odds': ' ',
                                  'Bad Rate': ' ',
                                  'KS': (np.sum(KS)).round(2),
                                  'Gini': (np.sum(gini)).round(2)})

        # Append the total row to ksGini DataFrame 
        ksGini = ksGini.append(total_row, ignore_index=True)

        # Calculate the total KS value  
        ksTotal = (np.sum(KS)).round(2)

        # Calculate the total Gini Value  
        giniTotal = (np.sum(gini)).round(2) 

        # Classify KS  and Gini  
        ks_summary = ks_classify(ksTotal, AppMon)
        gini_summary = gini_classify(giniTotal, AppMon)

    # POWER POINT 

    # TABLE SETTINGS: 
    Width1 = 1  
    Width2 = 1.5  
    Height = 1   
    yHeight = 0.25   
    Font = 10   

    pptx(file_name)  
    
    # TABLES 

    # TABLE KS AND GINI 
    table = tableDF(ksGini, 0.4, 3.4, 0.7, 0.7, 3.65, "Predictive Power - KS and Gini", "Personal Loan", " ", 'Monitoring', 10, 10, 10)  
    apply_formatting(table)   

    # TABLE KS CLASSIFICATION TABLE 
    table = tableDF(ks_summary, 8, yHeight, Width1, Width2, Height, "Predictive Power - KS and Gini", "Personal Loan", " ", 'KS', Font, Font, Font) 
    if not PerMon.empty:
        apply_formatting2(table, 'ks', ksTotal)
    else: 
        pass  

    # TABLE GINI CLASSIFICATION TABLE 
    table = tableDF(gini_summary, 10.7, yHeight, Width1, Width2, Height, "Predictive Power - KS and Gini", "Personal Loan", " ", 'Gini', Font, Font, Font) 
    if not PerMon.empty:
        apply_formatting2(table, 'gini', giniTotal)
    else:
        pass 
    
    # VISZUALTIONS 

    if not PerMon.empty:
        visualTwo("lineTwo", ksGini[:-1], ksGini[:-1], "Score Band", "Score Band", "No. of Bads","No. of Goods", "Distribution of Defaulters vs. NonDefaulters", "Score Band", "Distribution", 'Distribution of Defaulters', 'Distribution of Non-Defaulters', 1.57)  
    else: 
        visualThree("Distribution of Defaulters vs. NonDefaulters",'ScoreBand','Distribution', 1.57)  
        
    if not PerMon.empty:
        visualTwo("barLine", ksGini[:-1], ksGini[:-1], "Score Band","Score Band", "No. of Bads", "Bad Rate", "Bad Rate vs. Distribution", "Score Band", "Distribution of Bads",'Distribution of Bads', 'Bad Rate' , 4.38)     
    else: 
        visualThree('Bad Rate vs Distribution','ScoreBand','Distribution of Bads', 4.38)
    save(file_name) 
        
######################################
### THIRD SLIDE: CHARACTER ANALYSIS### 
######################################
    
    var_columns = AppDev.columns[AppDev.columns.str.startswith("var")]
    
    # Create a dictionary to store the dataframes
    tablesDev = {} 
    tablesMon = {} 
    # Development 
    
    # If App is empty
    if AppDev.empty:
        # Get the sheet names
        xls = pd.ExcelFile(file)
        sheet_names = xls.sheet_names[7:]  # Exclude the first 7 sheets

        # Iterate over each sheet and create a dataframe
        for sheet_name in sheet_names: 
            df = pd.read_excel(file, sheet_name=sheet_name) 
            tablesDev[sheet_name] = df 
            
    #If Per is empty 
    elif PerDev.empty: 
        for column in var_columns: 
            
            sheet_name = name_dict.get(column)
            
            # Application Characterstic Information 
            dev = calculate_char(AppDev, column)                                
            
            # Peformance Information 
            dev['No. of Bads'] = ''
            dev['No. of Goods'] = ''
            dev['Bad Rate'] = ''
            dev['WOE'] = ''
            dev['IV'] = ''
            
            dev.rename(columns={column: name_dict.get(column)}, inplace=True)
                

            # Total Row 
            totaldev = pd.DataFrame({name_dict.get(column): ['Total'],
                                  'Score': ' ',
                                  'Volume': np.sum(dev['Volume']),
                                  'Dist.': np.sum(dev['Dist.']).round(2),
                                  'No. of Bads': " " ,
                                  'No. of Goods': " ",
                                  'Bad Rate': ' ',
                                  'WOE': " ",
                                  'IV': " ", 
                                  })
            dev = dev.append(totaldev, ignore_index=True)
                             
            tablesDev[sheet_name] = dev
            
    # If App and Per is not empty         
    else: 
        for column in var_columns: 
            sheet_name = name_dict.get(column)
            
            # Application Characterstic Information 
            devGrouped = calculate_char(AppDev, column)                               
  

            # Performance Information 

            # Create Good and Bad Table 
            goodbadDev = calculate_perf(AppDev, PerDev, months, delinquency) 
            

            # Performance Characterstic Information 
            GBdev = calculate_char_two(goodbadDev, column)  
           

            # Concat Application Characterstic with Performance Information 
            dev = pd.merge(devGrouped, GBdev, on = column, how='inner')
         

            dev.rename(columns={column: name_dict.get(column)}, inplace=True)
            
            # Add Total Rows  
            totaldev = pd.DataFrame({name_dict.get(column): ['Total'],
                                  'Score': ' ',
                                  'Volume': np.sum(dev['Volume']),
                                  'Dist.': np.sum(dev['Dist.']).round(2),
                                  'No. of Bads': np.sum(dev['No. of Bads']),
                                  'No. of Goods': np.sum(dev['No. of Goods']),
                                  'Bad Rate': ' ',
                                  'WOE': np.sum(dev['WOE']).round(2),
                                  'IV':  np.sum(dev['IV']).round(2), 
                                  })
            dev = dev.append(totaldev, ignore_index=True)
                             
            tablesDev[sheet_name] = dev

    # Monitoring 
    
    # if App is empty  
    if AppMon.empty: 
        # Get the sheet names 
        xls = pd.ExcelFile(file)
        sheet_names = xls.sheet_names[7:]  # Exclude the first 7 sheets
        
        # Iterate over each sheet and create a dataframe
        for sheet_name in sheet_names:
            df = pd.read_excel(file, sheet_name=sheet_name)
            df['CSI'] = " "  
            tablesMon[sheet_name] = df  
            
        #CSI 
        CSI = 0 
        csi_summary = csi_classify(CSI, AppMon) 
        
        IV = 0 
        #IV 
        iv_summary = iv_classify(IV, AppMon)  
            
    #If Per is empty 
    elif PerMon.empty: 
        for column in var_columns: 
            sheet_name = name_dict.get(column)
            
            # Application Characterstic Information 
            mon = calculate_char(AppMon, column)                                
            
            # Peformance Information 
            mon['No. of Bads'] = ''
            mon['No. of Goods'] = ''
            mon['Bad Rate'] = ''
            mon['WOE'] = ''
            mon['IV'] = ''
            mon['CSI'] = ''
            
            mon.rename(columns={column: name_dict.get(column)}, inplace=True) 
                
            # Total Row 
            totalmon = pd.DataFrame({name_dict.get(column): ['Total'],
                                  'Score': ' ',
                                  'Volume': np.sum(mon['Volume']),
                                  'Dist.': np.sum(mon['Dist.']).round(2),
                                  'No. of Bads': " " ,
                                  'No. of Goods': " ",
                                  'Bad Rate': ' ',
                                  'WOE': " ",
                                  'IV': " ", 
                                  'CSI': " ", 
                                  })
            
            mon = mon.append(totalmon, ignore_index=True)
                             
            tablesMon[sheet_name] = mon
            
            CSI = 0 
            #CSI 
            csi_summary = csi_classify(CSI, PerMon) 
            
            IV = 0 
            #IV 
            iv_summary = iv_classify(IV, PerMon)  
            
    # If App and Per is not empty         
    else: 
        
        for column in var_columns: 
            # Application Characterstic Information 
            sheet_name = name_dict.get(column)
                                           
            monGrouped = calculate_char(AppMon, column)   

            # Performance Information 

            # Create Good and Bad Table 
            
            goodbadMon = calculate_perf(AppMon, PerMon, months, delinquency) 

            # Performance Characterstic Information 
            GBmon = calculate_char_two(goodBadMon, column) 

            # Concat Application Characterstic with Performance Information 
            mon = pd.merge(monGrouped, GBmon, on = column, how='inner')
            
            mon.rename(columns={column: name_dict.get(column)}, inplace=True)
                
            
            #CSI
            mon['CSI'] = ((mon['Volume'] - dev['Volume']) * np.log(mon['Volume'] / devGrouped['Volume'])).round(2)
            
            # Add Total Rows  
            totalmon = pd.DataFrame({name_dict.get(column): ['Total'],
                                  'Score': ' ',
                                  'Volume': np.sum(mon['Volume']),
                                  'Dist.': np.sum(mon['Dist.']).round(2),
                                  'No. of Bads': np.sum(mon['No. of Bads']),
                                  'No. of Goods': np.sum(mon['No. of Goods']),
                                  'Bad Rate': ' ',
                                  'WOE': np.sum(mon['WOE']).round(2),
                                  'IV':  np.sum(mon['IV']).round(2),
                                  'CSI':  np.sum(mon['CSI']).round(2)
                                  })
            mon = mon.append(totalmon, ignore_index=True) 
            
            tablesMon[sheet_name] = mon
    
            #CSI 
            csi_summary = csi_classify(np.sum(mon['CSI']).round(2), AppMon) 
    
            #IV 
            iv_summary = iv_classify(np.sum(mon['IV']).round(2), AppMon)  
            
############ IMPORT POWERPOINT TABLES ################# 
    
    for table in tablesDev.keys():              
    
        # TABLE SETTINGS:    
        Width1 = 1   
        Width2 = 1.5           
        Height = 1           
        yHeight = 0.25     
        Font = 10              

        dev = tablesDev[table]                     
        
        mon = tablesMon[table]    

        pptx(file_name)   

        #TABLE DEV 
        table = tableDF(dev, 0.15, 3, 1.5, 0.5, 1.8, "Character Analysis", "Personal Loan", " ", 'Development', 10, 10, 10)  
        apply_formatting(table) 

        #TABLE MON 
        table = tableDF(mon, 0.15, 5.12, 1.5, 0.5, 1.8, "Character Analysis", "Personal Loan", " ", 'Monitoring', 10, 10, 10) 
        apply_formatting(table) 

        # TABLE CSI CLASSIFICATION TABLE 
        table = tableDF(csi_summary, 8, yHeight, Width1, Width2, Height,"Character Analysis", "Personal Loan", " ", 'CSI', Font, Font, Font) 
        if not PerMon.empty: 
            apply_formatting2(table, 'csi', np.sum(mon['CSI']).round(2)) 
        else:  
            pass  

        # TABLE IV CLASSIFICATION TABLE 
        table = tableDF(iv_summary, 10.7, yHeight, Width1, Width2, Height, "Character Analysis", "Personal Loan", " ", 'IV', Font, Font, Font)  
        if not PerMon.empty: 
            apply_formatting2(table, 'iv', np.sum(mon['IV']).round(2))
        else: 
            pass  
        
        # VISUALIZATIONS 
        
        # CHARACTERISTIC ANALYSIS 
        if PerMon.empty:  
            visualThree('Bad Rate vs Distribution','Attribute','Distribution of Bads', 1.57)
        else: 
            pass  
            
        if not PerMon.empty:   
            visualTwo("barLine", mon.iloc[:-1], mon.iloc[:-1],  mon.columns[0], mon.columns[0], 'No. of Bads', 'Bad Rate', 'Distribution vs. Bad Rate', "Attribute", "Bads", "Monitoring Sample %", "Monitoring Bad Rate", 1.57)
        else: 
            pass 
        
        # WEIGHT OF EVIDENCE 
        if not PerDev.empty and not PerMon.empty:    
            visualTwo("lineTwo", dev.iloc[:-1], mon.iloc[:-1], dev.columns[0], dev[dev.columns[0]], "WOE", "WOE", 'Dev vs. Mon WOE', "Attribute", "WOE", 'Development WOE','Monitoring WOE' , 4.38)
        else: 
            pass 
        
        if PerDev.empty and not PerMon.empty: 
            visual("line", mon.iloc[:-1], mon.columns[0], "WOE", "Attribute", "WOE", "Monitoring WOE", 4.38) 
        else: 
            pass 
        
        if not PerDev.empty and PerMon.empty: 
            visual("line", dev.iloc[:-1], dev.columns[0], "WOE", "Attribute", "WOE", "Development WOE", 4.38)
        else: 
            pass 
        
        if PerDev.empty and PerMon.empty:   
            visualThree('Dev vs. Mon WOE','Attribute','WOE', 4.38)
        else: 
            pass  
        
        save(file_name)  
   
    
    
   
